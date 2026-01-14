#!/usr/bin/env python3
"""
Script per convertire PRESENTAZIONE_TECNICA.md in PowerPoint.

Requisiti:
    pip install python-pptx markdown

Uso:
    python convert_to_ppt.py
"""

import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def parse_markdown_slides(md_file):
    """Parse markdown file e estrae slide."""
    content = md_file.read_text(encoding='utf-8')
    
    # Divide per "---" che separa le slide
    slides_raw = re.split(r'\n---\n', content)
    
    slides = []
    for slide_raw in slides_raw:
        if not slide_raw.strip():
            continue
        
        lines = slide_raw.strip().split('\n')
        if not lines:
            continue
        
        # Estrai titolo (prima riga con ## o #)
        title = ""
        content_lines = []
        
        for line in lines:
            if line.startswith('## '):
                title = line[3:].strip()
            elif line.startswith('# '):
                title = line[2:].strip()
            else:
                content_lines.append(line)
        
        if title or content_lines:
            slides.append({
                'title': title,
                'content': '\n'.join(content_lines)
            })
    
    return slides

def get_graphic_for_slide(slide_num, title, content, used_graphics):
    """Determina quale grafico inserire in base al numero slide e contenuto.
    
    Args:
        slide_num: Numero della slide
        title: Titolo della slide
        content: Contenuto della slide
        used_graphics: Set di grafici già utilizzati
    
    Returns:
        Path al grafico o None se non trovato o già usato
    """
    graphics_dir = Path('grafici_presentazione')
    
    # Mappatura precisa slide -> grafico (solo una volta per grafico)
    slide_graphics = {
        4: '01_architettura.png',      # Slide 4: Architettura ad Alto Livello
        6: '03_stack_tecnologico.png',  # Slide 6: Stack Tecnologico
        9: '05_confronto_tts.png',     # Slide 9: Text-to-Speech (TTS)
        10: '02_flusso_e2e.png',        # Slide 10: Flusso Completo End-to-End
        11: '08_sequenza_multiuser.png', # Slide 11: Multi-User Support
        13: '07_database_schema.png',   # Slide 13: Database Schema
        14: '04_deployment_docker.png', # Slide 14: Deployment - Docker Compose
        24: '06_metriche_performance.png' # Slide 24: Metriche e Performance
    }
    
    # Usa solo la mappatura precisa, non keyword matching per evitare duplicati
    graphic = slide_graphics.get(slide_num)
    
    if graphic:
        # Controlla se questo grafico è già stato usato
        if graphic in used_graphics:
            return None  # Non inserire duplicati
        
        graphic_path = graphics_dir / graphic
        if graphic_path.exists():
            used_graphics.add(graphic)  # Segna come usato
            return graphic_path
    
    return None

def create_presentation(slides, output_file):
    """Crea presentazione PowerPoint dalle slide."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Layout titolo e contenuto
    title_slide_layout = prs.slide_layouts[0]  # Titolo
    content_slide_layout = prs.slide_layouts[1]  # Titolo e contenuto
    blank_layout = prs.slide_layouts[6]  # Layout vuoto per slide con immagini
    
    graphics_dir = Path('grafici_presentazione')
    has_graphics = graphics_dir.exists() and any(graphics_dir.glob('*.png'))
    used_graphics = set()  # Traccia grafici già utilizzati per evitare duplicati
    
    for i, slide_data in enumerate(slides):
        title = slide_data['title']
        content = slide_data['content']
        slide_num = i + 1
        
        # Controlla se questa slide ha un grafico (solo se non già usato)
        graphic_path = get_graphic_for_slide(slide_num, title, content, used_graphics)
        has_graphic = graphic_path and graphic_path.exists()
        
        # Prima slide: layout titolo
        if i == 0:
            slide = prs.slides.add_slide(title_slide_layout)
            title_shape = slide.shapes.title
            subtitle_shape = slide.placeholders[1]
            
            title_shape.text = title if title else "SophyAI Live Server"
            subtitle_shape.text = "Sistema di Assistente Vocale Self-Hosted"
        elif has_graphic:
            # Slide con grafico: usa layout con titolo e contenuto, aggiungi grafico e testo
            slide = prs.slides.add_slide(content_slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = title if title else f"Slide {slide_num}"
            
            # Aggiungi contenuto testuale sopra il grafico
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.word_wrap = True
            
            # Pulisci markdown e aggiungi testo introduttivo
            content_clean = clean_markdown(content)
            # Prendi solo le prime 2-3 righe per non occupare troppo spazio
            paragraphs = content_clean.split('\n\n')[:2]
            
            for para_text in paragraphs:
                if not para_text.strip():
                    continue
                if text_frame.paragraphs:
                    p = text_frame.add_paragraph()
                else:
                    p = text_frame.paragraphs[0]
                p.text = para_text.strip()[:200]  # Limita lunghezza
                p.level = 0
                p.font.size = Pt(12)
                p.space_after = Pt(4)
            
            # Aggiungi grafico sotto il testo (ridotto per lasciare spazio al testo)
            img_left = Inches(0.5)
            img_top = Inches(2.5)  # Più in basso per lasciare spazio al testo
            img_width = Inches(9)
            img_height = Inches(4.5)  # Leggermente più piccolo
            slide.shapes.add_picture(str(graphic_path), img_left, img_top, img_width, img_height)
            
            print(f"   📊 Slide {slide_num}: Inserito {graphic_path.name} (con testo)")
        else:
            # Altre slide: layout titolo e contenuto
            slide = prs.slides.add_slide(content_slide_layout)
            title_shape = slide.shapes.title
            content_shape = slide.placeholders[1]
            
            title_shape.text = title if title else f"Slide {slide_num}"
            
            # Formatta contenuto
            text_frame = content_shape.text_frame
            text_frame.word_wrap = True
            
            # Pulisci markdown e aggiungi testo
            content_clean = clean_markdown(content)
            paragraphs = content_clean.split('\n\n')
            
            for para_text in paragraphs:
                if not para_text.strip():
                    continue
                
                if text_frame.paragraphs:
                    p = text_frame.add_paragraph()
                else:
                    p = text_frame.paragraphs[0]
                
                p.text = para_text.strip()
                p.level = 0
                p.font.size = Pt(14)
                p.space_after = Pt(6)
    
    prs.save(output_file)
    print(f"✅ Presentazione creata: {output_file}")
    if has_graphics:
        print(f"📊 Grafici inseriti automaticamente nelle slide corrispondenti")

def clean_markdown(text):
    """Pulisce markdown da formattazione."""
    # Rimuovi code blocks
    text = re.sub(r'```[\s\S]*?```', '[CODE BLOCK]', text)
    # Rimuovi inline code
    text = re.sub(r'`([^`]+)`', r'\1', text)
    # Rimuovi link markdown
    text = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', text)
    # Rimuovi bold/italic
    text = re.sub(r'\*\*([^\*]+)\*\*', r'\1', text)
    text = re.sub(r'\*([^\*]+)\*', r'\1', text)
    # Rimuovi header markers
    text = re.sub(r'^###?\s+', '', text, flags=re.MULTILINE)
    # Rimuovi liste markdown
    text = re.sub(r'^\s*[-*+]\s+', '• ', text, flags=re.MULTILINE)
    text = re.sub(r'^\s*\d+\.\s+', '', text, flags=re.MULTILINE)
    
    return text

def main():
    md_file = Path('PRESENTAZIONE_TECNICA.md')
    output_file = Path('PRESENTAZIONE_TECNICA.pptx')
    
    if not md_file.exists():
        print(f"❌ File non trovato: {md_file}")
        return
    
    print(f"📖 Leggendo {md_file}...")
    slides = parse_markdown_slides(md_file)
    print(f"📊 Trovate {len(slides)} slide")
    
    print(f"📝 Creando PowerPoint...")
    create_presentation(slides, output_file)
    
    print(f"\n✅ Completato!")
    print(f"📄 File creato: {output_file}")
    print(f"\n💡 Suggerimenti:")
    print(f"   - Apri {output_file} in PowerPoint")
    print(f"   - Personalizza colori e stili")
    print(f"   - Aggiungi immagini/diagrammi dove necessario")

if __name__ == '__main__':
    try:
        from pptx import Presentation
    except ImportError:
        print("❌ python-pptx non installato!")
        print("   Installa con: pip install python-pptx")
        exit(1)
    
    main()
